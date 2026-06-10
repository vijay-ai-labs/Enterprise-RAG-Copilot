"""
Document ingestion: extract (text, page_number) tuples from PDF, TXT, MD, DOCX, PPTX, and HTML.
All formats produce page-wise or section-wise tuples so chunk metadata carries accurate page numbers.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import List, Tuple

from app.logger import logger

PageText = Tuple[str, int]  # (text, 1-indexed page/section number)


def extract_text(file_bytes: bytes, filename: str) -> List[PageText]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(file_bytes, filename)
    elif suffix in {".txt", ".md", ".markdown"}:
        return _extract_plaintext(file_bytes, filename)
    elif suffix == ".docx":
        return _extract_docx(file_bytes, filename)
    elif suffix == ".pptx":
        return _extract_pptx(file_bytes, filename)
    elif suffix in {".html", ".htm"}:
        return _extract_html(file_bytes, filename)
    else:
        raise ValueError(
            f"Unsupported file type: {suffix!r}. "
            "Supported: .pdf .txt .md .docx .pptx .html .htm"
        )


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _extract_pdf(file_bytes: bytes, filename: str) -> List[PageText]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf required: pip install pypdf") from exc

    reader = PdfReader(io.BytesIO(file_bytes))
    pages: List[PageText] = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append((text, page_num))

    logger.info("PDF '%s': extracted %d non-empty pages", filename, len(pages))
    return pages


# ---------------------------------------------------------------------------
# Plain text / Markdown
# ---------------------------------------------------------------------------

def _extract_plaintext(file_bytes: bytes, filename: str) -> List[PageText]:
    text = file_bytes.decode("utf-8", errors="replace").strip()
    logger.info("Plaintext '%s': %d chars", filename, len(text))
    return [(text, 1)]


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _extract_docx(file_bytes: bytes, filename: str) -> List[PageText]:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise RuntimeError("python-docx required: pip install python-docx") from exc

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
    except Exception as exc:
        raise RuntimeError(f"Failed to parse DOCX '{filename}': {exc}") from exc
    # Group paragraphs into ~page-sized sections (500 words ≈ 1 page)
    pages: List[PageText] = []
    current: List[str] = []
    word_count = 0
    page_num = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        current.append(text)
        word_count += len(text.split())
        if word_count >= 500:
            pages.append(("\n".join(current), page_num))
            current = []
            word_count = 0
            page_num += 1

    if current:
        pages.append(("\n".join(current), page_num))

    logger.info("DOCX '%s': %d sections extracted", filename, len(pages))
    return pages or [("", 1)]


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _extract_pptx(file_bytes: bytes, filename: str) -> List[PageText]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx required: pip install python-pptx") from exc

    prs = Presentation(io.BytesIO(file_bytes))
    pages: List[PageText] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append(("\n".join(texts), slide_num))

    logger.info("PPTX '%s': %d slides extracted", filename, len(pages))
    return pages


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def _extract_html(file_bytes: bytes, filename: str) -> List[PageText]:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise RuntimeError("beautifulsoup4 required: pip install beautifulsoup4 lxml") from exc

    soup = BeautifulSoup(file_bytes.decode("utf-8", errors="replace"), "lxml")

    # Remove boilerplate tags
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    text = "\n".join(line for line in text.splitlines() if line.strip())

    logger.info("HTML '%s': %d chars extracted", filename, len(text))
    return [(text, 1)] if text else []
